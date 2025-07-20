import streamlit as st
from io import BytesIO
import pandas as pd
from datetime import datetime
import numpy as np

# Import libraries for document generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    st.error("Please install python-pptx: pip install python-pptx")

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.units import inch
except ImportError:
    st.error("Please install reportlab: pip install reportlab")

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import LineChart, Reference
except ImportError:
    st.error("Please install openpyxl: pip install openpyxl")

try:
    from docx import Document
    from docx.shared import Inches as DocxInches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import RGBColor as DocxRGBColor
except ImportError:
    st.error("Please install python-docx: pip install python-docx")

class ProfessionalExporter:
    """Professional document export suite for Capital Markets Intelligence Platform"""
    
    def __init__(self, company_data, analytics_data, ticker):
        self.company_data = company_data
        self.analytics_data = analytics_data
        self.ticker = ticker
        self.generation_date = datetime.now().strftime("%B %d, %Y")
    
    def export_to_powerpoint(self):
        """Generate Goldman Sachs-style PowerPoint presentation"""
        try:
            prs = Presentation()
            
            # Slide 1: Title Slide
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = f"{self.company_data.get('name', self.ticker)} Investment Analysis"
            subtitle.text = f"Capital Markets Intelligence Platform\nGenerated on {self.generation_date}"
            
            # Slide 2: Executive Summary
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = "Executive Summary"
            
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = f"Security Analysis: {self.ticker}"
            
            # Add key metrics
            if self.analytics_data:
                p = tf.add_paragraph()
                p.text = f"• Expected Annual Return: {self.analytics_data.get('mean_return', 0)*100:.2f}%"
                p = tf.add_paragraph()
                p.text = f"• Risk-Adjusted Return (Sharpe): {self.analytics_data.get('sharpe_ratio', 0):.2f}"
                p = tf.add_paragraph()
                p.text = f"• Annualized Volatility: {self.analytics_data.get('volatility', 0)*100:.2f}%"
                p = tf.add_paragraph()
                p.text = f"• Maximum Drawdown: {self.analytics_data.get('max_drawdown', 0)*100:.2f}%"
            
            # Slide 3: Financial Metrics
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Key Financial Metrics"
            
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = f"Market Valuation - {self.ticker}"
            
            p = tf.add_paragraph()
            p.text = f"• Current Price: ${self.company_data.get('price', 0):.2f}"
            p = tf.add_paragraph()
            p.text = f"• Market Capitalization: ${self.company_data.get('market_cap', 0)/1e9:.1f}B"
            p = tf.add_paragraph()
            p.text = f"• P/E Ratio: {self.company_data.get('pe_ratio', 'N/A')}"
            p = tf.add_paragraph()
            p.text = f"• Trading Volume: {self.company_data.get('volume', 0):,}"
            
            # Save to BytesIO
            pptx_buffer = BytesIO()
            prs.save(pptx_buffer)
            pptx_buffer.seek(0)
            
            return pptx_buffer
            
        except Exception as e:
            st.error(f"PowerPoint generation error: {str(e)}")
            return None
    
    def export_to_pdf(self):
        """Generate professional PDF report"""
        try:
            buffer = BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4)
            
            # Get styles
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=colors.HexColor('#0f172a'),
                spaceAfter=30,
                alignment=1  # Center alignment
            )
            
            # Build document content
            story = []
            
            # Title
            story.append(Paragraph(f"{self.company_data.get('name', self.ticker)} Investment Analysis", title_style))
            story.append(Spacer(1, 20))
            
            # Executive Summary
            story.append(Paragraph("Executive Summary", styles['Heading2']))
            story.append(Paragraph(f"Comprehensive analysis of {self.ticker} conducted on {self.generation_date}", styles['Normal']))
            story.append(Spacer(1, 20))
            
            # Key Metrics Table
            if self.analytics_data:
                story.append(Paragraph("Key Financial Metrics", styles['Heading2']))
                
                data = [
                    ['Metric', 'Value'],
                    ['Expected Annual Return', f"{self.analytics_data.get('mean_return', 0)*100:.2f}%"],
                    ['Sharpe Ratio', f"{self.analytics_data.get('sharpe_ratio', 0):.2f}"],
                    ['Volatility', f"{self.analytics_data.get('volatility', 0)*100:.2f}%"],
                    ['Max Drawdown', f"{self.analytics_data.get('max_drawdown', 0)*100:.2f}%"],
                    ['VaR (95%)', f"{self.analytics_data.get('var_95', 0)*100:.2f}%"]
                ]
                
                table = Table(data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#0f172a')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f8fafc')),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                
                story.append(table)
            
            # Build PDF
            doc.build(story)
            buffer.seek(0)
            
            return buffer
            
        except Exception as e:
            st.error(f"PDF generation error: {str(e)}")
            return None
    
    def export_to_excel(self):
        """Generate comprehensive Excel financial model"""
        try:
            wb = Workbook()
            
            # Summary Sheet
            ws = wb.active
            ws.title = "Executive Summary"
            
            # Header styling
            header_font = Font(bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="0F172A", end_color="0F172A", fill_type="solid")
            
            # Title
            ws['A1'] = f"{self.company_data.get('name', self.ticker)} Investment Analysis"
            ws['A1'].font = Font(bold=True, size=16)
            ws.merge_cells('A1:D1')
            
            # Date
            ws['A2'] = f"Generated: {self.generation_date}"
            ws.merge_cells('A2:D2')
            
            # Headers
            ws['A4'] = "Metric"
            ws['B4'] = "Value"
            ws['C4'] = "Benchmark"
            ws['D4'] = "Rating"
            
            for cell in ['A4', 'B4', 'C4', 'D4']:
                ws[cell].font = header_font
                ws[cell].fill = header_fill
            
            # Data
            row = 5
            if self.analytics_data:
                metrics = [
                    ('Expected Annual Return', f"{self.analytics_data.get('mean_return', 0)*100:.2f}%", "10.0%", "Above Average"),
                    ('Sharpe Ratio', f"{self.analytics_data.get('sharpe_ratio', 0):.2f}", "1.0", "Good" if self.analytics_data.get('sharpe_ratio', 0) > 1 else "Below Average"),
                    ('Volatility', f"{self.analytics_data.get('volatility', 0)*100:.2f}%", "15.0%", "Normal"),
                    ('Max Drawdown', f"{self.analytics_data.get('max_drawdown', 0)*100:.2f}%", "-10.0%", "Concerning" if self.analytics_data.get('max_drawdown', 0) < -0.15 else "Acceptable")
                ]
                
                for metric, value, benchmark, rating in metrics:
                    ws[f'A{row}'] = metric
                    ws[f'B{row}'] = value
                    ws[f'C{row}'] = benchmark
                    ws[f'D{row}'] = rating
                    row += 1
            
            # Company Data Sheet
            ws2 = wb.create_sheet("Company Data")
            ws2['A1'] = "Company Information"
            ws2['A1'].font = Font(bold=True, size=14)
            
            company_info = [
                ('Ticker', self.ticker),
                ('Company Name', self.company_data.get('name', 'N/A')),
                ('Current Price', f"${self.company_data.get('price', 0):.2f}"),
                ('Market Cap', f"${self.company_data.get('market_cap', 0)/1e9:.1f}B"),
                ('P/E Ratio', str(self.company_data.get('pe_ratio', 'N/A'))),
                ('Sector', self.company_data.get('sector', 'N/A')),
                ('Industry', self.company_data.get('industry', 'N/A'))
            ]
            
            row = 3
            for label, value in company_info:
                ws2[f'A{row}'] = label
                ws2[f'B{row}'] = value
                row += 1
            
            # Auto-adjust column widths
            for ws in wb.worksheets:
                for column in ws.columns:
                    max_length = 0
                    column_letter = column[0].column_letter
                    for cell in column:
                        try:
                            if len(str(cell.value)) > max_length:
                                max_length = len(str(cell.value))
                        except:
                            pass
                    adjusted_width = min(max_length + 2, 50)
                    ws.column_dimensions[column_letter].width = adjusted_width
            
            # Save to BytesIO
            excel_buffer = BytesIO()
            wb.save(excel_buffer)
            excel_buffer.seek(0)
            
            return excel_buffer
            
        except Exception as e:
            st.error(f"Excel generation error: {str(e)}")
            return None
    
    def export_to_word(self):
        """Generate detailed Word analysis document"""
        try:
            doc = Document()
            
            # Title
            title = doc.add_heading(f"{self.company_data.get('name', self.ticker)} Investment Analysis", 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Date
            date_para = doc.add_paragraph(f"Generated on {self.generation_date}")
            date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_page_break()
            
            # Executive Summary
            doc.add_heading('Executive Summary', level=1)
            doc.add_paragraph(f"This report provides a comprehensive investment analysis of {self.ticker} "
                            f"({self.company_data.get('name', 'N/A')}) based on advanced quantitative methods "
                            "and market intelligence.")
            
            # Company Overview
            doc.add_heading('Company Overview', level=1)
            overview_table = doc.add_table(rows=1, cols=2)
            overview_table.style = 'Table Grid'
            
            hdr_cells = overview_table.rows[0].cells
            hdr_cells[0].text = 'Attribute'
            hdr_cells[1].text = 'Value'
            
            company_details = [
                ('Ticker Symbol', self.ticker),
                ('Company Name', self.company_data.get('name', 'N/A')),
                ('Current Price', f"${self.company_data.get('price', 0):.2f}"),
                ('Market Capitalization', f"${self.company_data.get('market_cap', 0)/1e9:.1f} Billion"),
                ('Sector', self.company_data.get('sector', 'N/A')),
                ('Industry', self.company_data.get('industry', 'N/A'))
            ]
            
            for attr, value in company_details:
                row_cells = overview_table.add_row().cells
                row_cells[0].text = attr
                row_cells[1].text = str(value)
            
            # Financial Analysis
            if self.analytics_data:
                doc.add_heading('Quantitative Analysis', level=1)
                
                analysis_text = (
                    f"Our quantitative analysis reveals the following key insights for {self.ticker}:\n\n"
                    f"• Expected Annual Return: {self.analytics_data.get('mean_return', 0)*100:.2f}%\n"
                    f"• Risk-Adjusted Return (Sharpe Ratio): {self.analytics_data.get('sharpe_ratio', 0):.2f}\n"
                    f"• Annualized Volatility: {self.analytics_data.get('volatility', 0)*100:.2f}%\n"
                    f"• Maximum Drawdown: {self.analytics_data.get('max_drawdown', 0)*100:.2f}%\n"
                    f"• Value at Risk (95% confidence): {self.analytics_data.get('var_95', 0)*100:.2f}%"
                )
                
                doc.add_paragraph(analysis_text)
            
            # Risk Assessment
            doc.add_heading('Risk Assessment', level=1)
            doc.add_paragraph("Based on our analysis, this security presents the following risk profile:")
            
            if self.analytics_data:
                volatility = self.analytics_data.get('volatility', 0)
                if volatility > 0.25:
                    risk_level = "High Risk"
                elif volatility > 0.15:
                    risk_level = "Moderate Risk"
                else:
                    risk_level = "Low Risk"
                
                doc.add_paragraph(f"Risk Classification: {risk_level}")
            
            # Save to BytesIO
            word_buffer = BytesIO()
            doc.save(word_buffer)
            word_buffer.seek(0)
            
            return word_buffer
            
        except Exception as e:
            st.error(f"Word document generation error: {str(e)}")
            return None

def get_export_suite():
    """Return the ProfessionalExporter class for use in main app"""
    return ProfessionalExporter
